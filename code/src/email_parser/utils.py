import base64
import pytesseract
from PIL import Image
import io
import docx
import fitz
from pptx import Presentation
import re
from src.models import Attachment
from typing import List, Optional, Dict


def decode_base64(base64_content):
    return base64.b64decode(base64_content)


def extract_text_from_pdf(file_data):
    text = ""
    with io.BytesIO(file_data) as pdf_file:
        doc = fitz.open(stream=pdf_file, filetype="pdf")
        for page in doc:
            text += page.get_text("text")
    text = re.sub(r"\s+", " ", text).strip()
    return text


def extract_text_from_docx(file_data):
    text = ""
    with io.BytesIO(file_data) as docx_file:
        doc = docx.Document(docx_file)
        for para in doc.paragraphs:
            text += para.text + "\n"
    text = re.sub(r"\s+", " ", text).strip()
    return text


def extract_text_from_pptx(file_data):
    text = ""
    with io.BytesIO(file_data) as pptx_file:
        presentation = Presentation(pptx_file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    text = re.sub(r"\s+", " ", text).strip()
    return text


def extract_text_from_image(file_data):
    image = Image.open(io.BytesIO(file_data))
    return pytesseract.image_to_string(image)


def process_attachment(filename, base64_content):
    file_data = decode_base64(base64_content)
    file_text = ""

    if filename.endswith(".pdf"):
        file_text = extract_text_from_pdf(file_data)
    elif filename.endswith(".docx"):
        file_text = extract_text_from_docx(file_data)
    elif filename.endswith(".pptx"):
        file_text = extract_text_from_pptx(file_data)
    elif (
        filename.endswith(".jpg")
        or filename.endswith(".png")
        or filename.endswith(".jpeg")
    ):
        file_text = extract_text_from_image(file_data)
    else:
        file_text = "Unsupported file type."

    return file_text


# Main function to parse attachments
def parse_attachments(attachments):
    parsed_attachments = []

    for attachment in attachments:
        filename = attachment.fileName
        content = attachment.data

        extracted_content = process_attachment(filename, content)

        parsed_attachments.append(Attachment(fileName=filename, data=extracted_content))

    return parsed_attachments
